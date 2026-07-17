"""Локальные парсеры: они не выполняют формулы, макросы или содержимое файлов."""

from __future__ import annotations

import csv
import json
import math
import shutil
import subprocess
import unicodedata
from datetime import date, datetime, time
from pathlib import Path
from typing import Any, NamedTuple
from xml.etree import ElementTree

from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from ..config import LimitsSettings


class ParsedText(NamedTuple):
    text: str
    warnings: list[str]
    page_count: int | None
    usable: bool


def _decode(path: Path) -> str:
    return path.read_bytes().decode("utf-8", errors="replace")


def _doc_converters() -> list[tuple[str, list[str]]]:
    """Возвращает только локальные конвертеры legacy DOC без запуска макросов."""

    converters: list[tuple[str, list[str]]] = []
    if Path("/usr/bin/textutil").is_file():
        converters.append(("/usr/bin/textutil", ["-convert", "txt", "-stdout"]))
    if antiword := shutil.which("antiword"):
        converters.append((antiword, ["-w", "0"]))
    return converters


def _doc_text(path: Path) -> tuple[str, list[str]]:
    """Извлекает текст DOC системным конвертером, не открывая документ в Office."""

    for executable, arguments in _doc_converters():
        try:
            result = subprocess.run(
                [executable, *arguments, str(path)],
                stdin=subprocess.DEVNULL,
                capture_output=True,
                check=False,
                timeout=30,
            )
        except (OSError, subprocess.TimeoutExpired):
            continue
        if result.returncode != 0:
            continue
        text = result.stdout.decode("utf-8", errors="replace").strip()
        if text:
            return text, []
    return "", ["Не удалось извлечь текст из DOC; требуется ручная проверка файла."]


def sanitize_html(value: str) -> str:
    soup = BeautifulSoup(value, "html.parser")
    for node in soup(["script", "style", "iframe", "object", "embed", "form", "base", "meta"]):
        node.decompose()
    for node in soup.find_all(style=True):
        style = str(node.get("style", "")).replace(" ", "").lower()
        if "display:none" in style or "visibility:hidden" in style:
            node.decompose()
    return soup.get_text("\n", strip=True)


def _limited(value: str, limit: int = 400_000) -> tuple[str, list[str]]:
    if len(value) <= limit:
        return value, []
    return value[:limit] + "\n[truncated by agent limit]", ["Локальное извлечение ограничено лимитом текста."]


def _cell_text(value: object) -> str:
    """Безопасно представляет значение ячейки, не вычисляя формулы."""

    if value is None:
        return ""
    if isinstance(value, (datetime, date, time)):
        return value.isoformat(sep=" ") if isinstance(value, datetime) else value.isoformat()
    if isinstance(value, float):
        if math.isnan(value) or math.isinf(value):
            return str(value)
        if value.is_integer():
            return str(int(value))
    text = " ".join(str(value).split())
    return "".join(character for character in text if character.isprintable())


def _spreadsheet_is_garbled(value: str) -> bool:
    """Отмечает только явно повреждённый текст, не пытаясь угадывать смысл ячеек."""

    replacement = value.count("�")
    controls = sum(unicodedata.category(character).startswith("C") and character not in "\n\r\t" for character in value)
    return replacement > 0 or controls > max(3, len(value) // 200)


def _is_textual(value: str) -> bool:
    try:
        float(value.replace(" ", "").replace(",", "."))
    except ValueError:
        return bool(value)
    return False


def _unique_headers(values: list[str]) -> list[str]:
    seen: dict[str, int] = {}
    headers: list[str] = []
    for index, value in enumerate(values, 1):
        base = value or f"Колонка {index}"
        count = seen.get(base, 0) + 1
        seen[base] = count
        headers.append(base if count == 1 else f"{base} ({count})")
    return headers


def _spreadsheet_sheet_text(title: str, rows: list[tuple[int, list[str]]]) -> list[str]:
    """Сериализует лист как подписи столбцов и строки, а не как неразмеченный TSV."""

    populated = [(number, cells) for number, cells in rows if any(cells)]
    width = max((len([cell for cell in cells if cell]) for _, cells in populated), default=0)
    lines = [f'[Лист: "{title}"; строк с данными: {len(populated)}; заполненных колонок: {width}]']
    if not populated:
        lines.append("[Лист пуст]")
        return lines

    candidates = populated[:30]
    header_width = max((len([cell for cell in cells if cell]) for _, cells in candidates), default=1)
    threshold = min(3, header_width)
    header_position = 0
    for position, (_, cells) in enumerate(candidates):
        nonempty = [cell for cell in cells if cell]
        if len(nonempty) >= threshold and any(_is_textual(cell) for cell in nonempty):
            header_position = position
            break
    header_number, header_values = populated[header_position]
    headers = _unique_headers(header_values)
    lines.append("[Заголовки: " + " | ".join(headers) + "]")

    for number, cells in populated[:header_position]:
        values = [value for value in cells if value]
        lines.append(f"Преамбула, строка {number}: " + " | ".join(values))
    for number, cells in populated[header_position + 1 :]:
        pairs = [f"{headers[index]}={value}" for index, value in enumerate(cells) if value]
        if not pairs:
            continue
        is_total = any(value.casefold().startswith(("итого", "всего", "total")) for value in cells if value)
        prefix = "Итог, " if is_total else ""
        lines.append(f"{prefix}строка {number}: " + "; ".join(pairs))
    return lines


def _xlsx_text(path: Path, limits: LimitsSettings) -> tuple[str, list[str]]:
    """Читает кэшированные результаты XLSX-формул без их исполнения."""

    values_book = load_workbook(path, read_only=True, data_only=True)
    formulas_book = load_workbook(path, read_only=True, data_only=False)
    warnings: list[str] = []
    values: list[str] = ["[Таблица: XLSX; формулы не выполнялись, использованы сохранённые результаты]"]
    has_values = False
    try:
        for worksheet in values_book.worksheets[: limits.max_xlsx_sheets]:
            formula_sheet = formulas_book[worksheet.title]
            rows: list[tuple[int, list[str]]] = []
            value_rows = worksheet.iter_rows(
                min_row=1,
                max_row=min(worksheet.max_row, limits.max_xlsx_rows),
                max_col=min(worksheet.max_column, limits.max_xlsx_columns),
                values_only=False,
            )
            formula_rows = formula_sheet.iter_rows(
                min_row=1,
                max_row=min(formula_sheet.max_row, limits.max_xlsx_rows),
                max_col=min(formula_sheet.max_column, limits.max_xlsx_columns),
                values_only=False,
            )
            for number, (value_row, formula_row) in enumerate(zip(value_rows, formula_rows, strict=True), 1):
                cells: list[str] = []
                for value_cell, formula_cell in zip(value_row, formula_row, strict=True):
                    value = _cell_text(value_cell.value)
                    formula = formula_cell.value
                    if isinstance(formula, str) and formula.startswith("="):
                        value = (
                            f"{value} [формула: {formula}]"
                            if value
                            else f"[формула без сохранённого результата: {formula}]"
                        )
                    cells.append(value)
                rows.append((number, cells))
            values.extend(_spreadsheet_sheet_text(worksheet.title, rows))
            has_values = has_values or any(any(cells) for _, cells in rows)
            if worksheet.max_row > limits.max_xlsx_rows:
                warnings.append(f"Лист {worksheet.title}: достигнут лимит строк.")
            if worksheet.max_column > limits.max_xlsx_columns:
                warnings.append(f"Лист {worksheet.title}: достигнут лимит колонок.")
        if len(values_book.worksheets) > limits.max_xlsx_sheets:
            warnings.append("Достигнут лимит листов XLSX.")
    finally:
        values_book.close()
        formulas_book.close()
    if not has_values:
        warnings.append("В XLSX нет читаемых значений ячеек.")
    return "\n".join(values), warnings


def _xls_cell_text(cell: Any, datemode: int, xlrd_module: Any) -> str:
    if cell.ctype == xlrd_module.XL_CELL_DATE:
        return _cell_text(xlrd_module.xldate_as_datetime(cell.value, datemode))
    return _cell_text(cell.value)


def _xls_text(path: Path, limits: LimitsSettings) -> tuple[str, list[str]]:
    """Читает legacy XLS через xlrd; формулы и макросы не исполняются."""

    import xlrd

    workbook = xlrd.open_workbook(path, on_demand=True)
    warnings: list[str] = []
    values: list[str] = ["[Таблица: XLS; использованы сохранённые значения ячеек]"]
    has_values = False
    try:
        for worksheet in workbook.sheets()[: limits.max_xlsx_sheets]:
            rows: list[tuple[int, list[str]]] = []
            for row_index in range(min(worksheet.nrows, limits.max_xlsx_rows)):
                cells = [
                    _xls_cell_text(worksheet.cell(row_index, column), workbook.datemode, xlrd)
                    for column in range(min(worksheet.ncols, limits.max_xlsx_columns))
                ]
                rows.append((row_index + 1, cells))
            values.extend(_spreadsheet_sheet_text(worksheet.name, rows))
            has_values = has_values or any(any(cells) for _, cells in rows)
            if worksheet.nrows > limits.max_xlsx_rows:
                warnings.append(f"Лист {worksheet.name}: достигнут лимит строк.")
            if worksheet.ncols > limits.max_xlsx_columns:
                warnings.append(f"Лист {worksheet.name}: достигнут лимит колонок.")
        if workbook.nsheets > limits.max_xlsx_sheets:
            warnings.append("Достигнут лимит листов XLS.")
    finally:
        workbook.release_resources()
    if not has_values:
        warnings.append("В XLS нет читаемых значений ячеек.")
    return "\n".join(values), warnings


def extract_programmatic(path: Path, extension: str, limits: LimitsSettings) -> ParsedText:
    """Возвращает текст и явный статус качества без выполнения пользовательского кода."""

    extension = extension.lower()
    warnings: list[str] = []
    if extension in {".txt", ".md"}:
        text = _decode(path)
    elif extension in {".html", ".htm"}:
        text = sanitize_html(_decode(path))
    elif extension == ".json":
        text = json.dumps(json.loads(_decode(path)), ensure_ascii=False, indent=2)
    elif extension == ".xml":
        root = ElementTree.fromstring(path.read_bytes())
        text = "\n".join(part.strip() for part in root.itertext() if part.strip())
    elif extension in {".csv", ".tsv"}:
        delimiter = "\t" if extension == ".tsv" else ","
        with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
            rows = list(csv.reader(handle, delimiter=delimiter))
        text = "\n".join("\t".join(row) for row in rows)
    elif extension == ".docx":
        document = Document(str(path))
        document_values = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        for table in document.tables:
            document_values.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
        text = "\n".join(document_values)
    elif extension == ".doc":
        text, doc_warnings = _doc_text(path)
        warnings.extend(doc_warnings)
    elif extension == ".xlsx":
        text, spreadsheet_warnings = _xlsx_text(path, limits)
        warnings.extend(spreadsheet_warnings)
    elif extension == ".xls":
        text, spreadsheet_warnings = _xls_text(path, limits)
        warnings.extend(spreadsheet_warnings)
    elif extension == ".pptx":
        presentation = Presentation(str(path))
        values = []
        for number, slide in enumerate(presentation.slides[: limits.max_pptx_slides], start=1):
            values.append(f"[slide: {number}]")
            values.extend(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text)
            notes = getattr(slide, "notes_slide", None)
            if notes is not None:
                values.extend(shape.text for shape in notes.shapes if hasattr(shape, "text") and shape.text)
        if len(presentation.slides) > limits.max_pptx_slides:
            warnings.append("Достигнут лимит слайдов PPTX.")
        text = "\n".join(values)
    elif extension == ".pdf":
        reader = PdfReader(path)
        pages = len(reader.pages)
        if pages > limits.max_pdf_pages:
            warnings.append("Достигнут лимит страниц PDF.")
        text = "\n".join((page.extract_text() or "") for page in reader.pages[: limits.max_pdf_pages])
        text, truncation = _limited(text)
        warnings.extend(truncation)
        usable = len(text.strip()) >= max(20, min(200, pages * 10)) and text.count("�") < max(3, len(text) // 100)
        return ParsedText(text, warnings, pages, usable)
    else:
        return ParsedText("", ["Формат не поддерживается локальным парсером."], None, False)
    text, truncation = _limited(text)
    warnings.extend(truncation)
    spreadsheet_invalid = extension in {".xls", ".xlsx"} and (
        _spreadsheet_is_garbled(text) or any("нет читаемых значений" in warning for warning in warnings)
    )
    if spreadsheet_invalid:
        warnings.append("Не удалось корректно извлечь содержимое таблицы; требуется ручная проверка файла.")
    return ParsedText(text, warnings, None, bool(text.strip()) and not spreadsheet_invalid)
